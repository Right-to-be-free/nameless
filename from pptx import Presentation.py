from pptx import Presentation

# Create a new PowerPoint presentation
prs = Presentation()

def add_slide(prs, title, content):
    """Adds a slide with a title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    textbox = slide.placeholders[1]
    textbox.text = content

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Understanding YouTube Data Through Mathematical Concepts"
slide.placeholders[1].text = "An Analytical Approach to Video Trends"

# Content for slides
slides_content = [
    ("Introduction", "Purpose of the analysis\nImportance of understanding data trends\nHow mathematical concepts help analyze YouTube data"),
    ("Scalars", "Definition: A single numerical value\nExample: Number of views, likes, dislikes per video"),
    ("Vectors", "Definition: An ordered list of numbers\nExample: A list of likes and dislikes across multiple videos"),
    ("Matrices", "Definition: A two-dimensional array of numbers\nExample: Dataset table with views, likes, dislikes"),
    ("Basis", "Definition: Fundamental components that span a space\nExample: Views, likes, dislikes forming the basis of analysis"),
    ("Eigenvectors & Eigenvalues", "Definition: Special vectors that remain unchanged after transformation\nExample: PCA analysis in YouTube trends"),
    ("Mean, Median, Mode", "Mean: Average views\nMedian: Middle value in sorted views\nMode: Most common video category"),
    ("Bernoulli’s Distribution", "Definition: Binary outcome probability\nExample: 'Viral or Not?' classification based on 1M+ views"),
    ("Normal Distribution", "Definition: Bell curve where most data points lie near the mean\nExample: Video view distribution across all categories"),
    ("Key Takeaways", "Summary of insights\nHow businesses can leverage YouTube trends\nNext steps for decision-making"),
    ("Q&A", "Open floor for discussion\nContact information"),
]

# Add slides dynamically
for title, content in slides_content:
    add_slide(prs, title, content)

# Save the presentation
file_name = "YouTube_Data_Analysis_Presentation.pptx"
prs.save(file_name)
print(f"Presentation '{file_name}' created successfully!")
