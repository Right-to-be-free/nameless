from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Python Object-Oriented Programming (OOPs)"
subtitle.text = "A detailed explanation of OOPs concepts in Python"

# Slide 1: Introduction to OOPs
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Introduction to OOPs"
content.text = "• Object-Oriented Programming (OOP) is a paradigm based on objects and classes.\n" \
               "• It allows modular, reusable, and scalable code design.\n" \
               "• Key concepts include Class, Object, Attributes, Methods, Inheritance, Encapsulation, and Polymorphism."

# Slide 2: Class and Object
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Class and Object"
content.text = "• Class: A blueprint for creating objects.\n" \
               "• Object: An instance of a class.\n" \
               "Example:\nclass Car:\n    def __init__(self, brand):\n        self.brand = brand\ncar1 = Car('Toyota')"

# Slide 3: Attributes and Methods
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Attributes and Methods"
content.text = "• Attributes: Variables that store object state.\n" \
               "• Methods: Functions inside a class that define behavior.\n" \
               "Example:\nclass Car:\n    def __init__(self, brand):\n        self.brand = brand\n" \
               "    def show_brand(self):\n        return self.brand\ncar1.show_brand()"

# Slide 4: Constructors and Destructors
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Constructors and Destructors"
content.text = "• Constructor (__init__): Initializes object properties.\n" \
               "• Destructor (__del__): Cleans up resources.\n" \
               "Example:\nclass Car:\n    def __init__(self, brand):\n        self.brand = brand\n" \
               "    def __del__(self):\n        print('Object deleted')"

# Slide 5: Inheritance
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Inheritance"
content.text = "• Inheritance allows a class to acquire properties of another class.\n" \
               "• Supports code reusability.\n" \
               "Example:\nclass Car:\n    def __init__(self, brand):\n        self.brand = brand\n" \
               "class ElectricCar(Car):\n    def __init__(self, brand, battery):\n        super().__init__(brand)\n        self.battery = battery"

# Slide 6: Encapsulation
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Encapsulation"
content.text = "• Hides object details from the outside world.\n" \
               "• Private attributes start with '__'.\n" \
               "Example:\nclass Car:\n    def __init__(self, brand):\n        self.__brand = brand\n    def get_brand(self):\n        return self.__brand"

# Slide 7: Polymorphism
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Polymorphism"
content.text = "• Same function name, different implementations.\n" \
               "• Example: Method Overriding.\n" \
               "Example:\nclass Car:\n    def drive(self):\n        return 'Driving a car'\nclass Bike:\n    def drive(self):\n        return 'Riding a bike'\nfor v in [Car(), Bike()]:\n    print(v.drive())"

# Slide 8: *args and **kwargs
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "*args and **kwargs"
content.text = "• *args: Passes a variable number of arguments.\n" \
               "• **kwargs: Passes named arguments as a dictionary.\n" \
               "Example:\ndef example(*args, **kwargs):\n    print(args)\n    print(kwargs)\nexample(1, 2, name='John', age=30)"

# Save the presentation
pptx_path = "C:/Users/rishi/Desktop/Gen Ai/Python_OOPs_Presentation.pptx"
prs.save(pptx_path)

print(f"Presentation saved at: {pptx_path}")